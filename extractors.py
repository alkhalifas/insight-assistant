# extractors.py
import pandas as pd
from pptx import Presentation

def extract_text_from_pptx(file_path):
    prs = Presentation(file_path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

def load_csv(file_path):
    return pd.read_csv(file_path)
